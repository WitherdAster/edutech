"""
EduTech - Generate PPTX Presentation
Proposal Tugas Akhir: Face Recognition Attendance System
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from PIL import Image as PILImage
import os

# ─── CONSTANTS ───────────────────────────────────────────────────
BLUE_PRIMARY = RGBColor(0x4A, 0x6C, 0xF7)
BLUE_DARK = RGBColor(0x35, 0x58, 0xE0)
BLUE_LIGHT = RGBColor(0xE8, 0xED, 0xFF)
CYAN = RGBColor(0x0E, 0xA5, 0xE9)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x1A, 0x1A, 0x2E)
GRAY = RGBColor(0x6B, 0x72, 0x80)
GRAY_LIGHT = RGBColor(0xF5, 0xF7, 0xFA)
GREEN = RGBColor(0x10, 0xB9, 0x81)

LOGO_PATH = r'C:\Users\advan\Documents\EduTech\docs\polihasnurlogo.png'
OUTPUT_PATH = r'C:\Users\advan\Documents\EduTech\docs\presentasi_proposal.pptx'

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
CONTENT_MARGIN_L = Inches(1.2)
CONTENT_MARGIN_R = Inches(1.2)
CONTENT_TOP = Inches(1.6)
CONTENT_W = SLIDE_W - CONTENT_MARGIN_L - CONTENT_MARGIN_R

# ─── HELPERS ─────────────────────────────────────────────────────

def set_morph_transition(slide):
    """Set Morph transition on a slide via XML manipulation."""
    sld = slide._element
    # Remove existing transition if any
    existing = sld.find(qn('p:transition'))
    if existing is not None:
        sld.remove(existing)
    # Create transition element
    trans = sld.makeelement(qn('p:transition'), {})
    trans.set(qn('p:advanceOnClick'), '1')
    trans.set(qn('p:advanceMode'), 'onClick')
    # Add morph child
    morph = trans.makeelement(qn('p:morph'), {})
    trans.append(morph)
    # Insert as first child
    sld.insert(0, trans)


def add_bg_rect(slide, left, top, width, height, color):
    """Add a solid rectangle as background element."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_textbox(slide, left, top, width, height, text, font_size=18,
                bold=False, color=BLACK, alignment=PP_ALIGN.LEFT,
                font_name='Calibri', anchor=MSO_ANCHOR.TOP, spacing=1.2):
    """Add a text box with formatted text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    # Set anchor
    tf.paragraphs[0].alignment = alignment
    # Set vertical anchor
    txBox.text_frame._txBody.bodyPr.set('anchor', {
        MSO_ANCHOR.TOP: 't',
        MSO_ANCHOR.MIDDLE: 'ctr',
        MSO_ANCHOR.BOTTOM: 'b',
    }.get(anchor, 't'))

    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    p.space_after = Pt(0)
    p.space_before = Pt(0)
    return txBox


def add_bullet_text(slide, left, top, width, height, items,
                    font_size=16, color=BLACK, font_name='Calibri',
                    bullet_char='\u2022', line_spacing=1.5):
    """Add a text box with bullet points."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f'{bullet_char} {item}'
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.space_after = Pt(6)
        p.space_before = Pt(2)
        p.line_spacing = line_spacing
    return txBox


def add_slide_number(slide, number, total=10):
    """Add slide number at bottom right."""
    add_textbox(slide, Inches(11.5), Inches(7.0), Inches(1.5), Inches(0.4),
                f'{number} / {total}', font_size=11, color=GRAY,
                alignment=PP_ALIGN.RIGHT)


def add_section_header(slide, title, subtitle=None, slide_num=None, total=10):
    """Add a consistent blue header bar to a slide."""
    # Blue header band
    add_bg_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.2), BLUE_PRIMARY)
    # Title text
    add_textbox(slide, CONTENT_MARGIN_L, Inches(0.2), CONTENT_W, Inches(0.8),
                title, font_size=28, bold=True, color=WHITE, font_name='Calibri')
    # Subtitle if provided
    if subtitle:
        add_textbox(slide, CONTENT_MARGIN_L, Inches(0.75), CONTENT_W, Inches(0.4),
                    subtitle, font_size=14, color=RGBColor(0xCC, 0xD5, 0xFF),
                    font_name='Calibri')
    # Thin accent line below header
    add_bg_rect(slide, Inches(0), Inches(1.2), SLIDE_W, Inches(0.04), CYAN)
    if slide_num:
        add_slide_number(slide, slide_num, total)


def add_footer_bar(slide):
    """Add a thin gray line at the bottom."""
    add_bg_rect(slide, Inches(0), Inches(7.35), SLIDE_W, Inches(0.02), GRAY)


# ─── SLIDE BUILDERS ──────────────────────────────────────────────

def create_slide1(pres):
    """Slide 1: Cover"""
    slide = pres.slides.add_slide(pres.slide_layouts[6])  # blank

    # Full blue background
    add_bg_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(4.5), BLUE_PRIMARY)
    # White bottom half
    add_bg_rect(slide, Inches(0), Inches(4.5), SLIDE_W, Inches(3.0), WHITE)

    # Logo
    if os.path.exists(LOGO_PATH):
        # Resize logo to fit nicely
        img = PILImage.open(LOGO_PATH)
        w, h = img.size
        # Target width ~1.5 inches
        scale = 1.5 / (w / 96)
        logo_w = Inches(1.5)
        logo_h = Inches(h / 96 * scale)
        logo_left = (SLIDE_W - logo_w) // 2
        slide.shapes.add_picture(LOGO_PATH, logo_left, Inches(0.6), logo_w, logo_h)

    # Title
    title_text = 'RANCANG BANGUN SISTEM PRESENSI SISWA\nDENGAN FITUR FACE RECOGNITION\nBERBASIS WEB DAN ANDROID\nDI SMKN 1 MARABAHAN'
    add_textbox(slide, Inches(1.5), Inches(2.2), Inches(10.3), Inches(2.0),
                title_text, font_size=24, bold=True, color=WHITE,
                alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                font_name='Calibri')

    # Proposal tag
    add_textbox(slide, Inches(4), Inches(4.7), Inches(5.3), Inches(0.5),
                'PROPOSAL TUGAS AKHIR', font_size=14, bold=True, color=BLUE_PRIMARY,
                alignment=PP_ALIGN.CENTER, font_name='Calibri')

    # Author info
    add_textbox(slide, Inches(2), Inches(5.3), Inches(9.3), Inches(1.5),
                'Ahmad Nuur Fu\'ady  |  24302064\n'
                'Program Studi D3 Teknik Informatika\n'
                'Politeknik Hasnur',
                font_size=16, color=GRAY, alignment=PP_ALIGN.CENTER,
                font_name='Calibri')

    add_slide_number(slide, 1)
    set_morph_transition(slide)
    return slide


def create_slide2(pres):
    """Slide 2: Latar Belakang Masalah"""
    slide = pres.slides.add_slide(pres.slide_layouts[6])
    add_section_header(slide, 'Latar Belakang Masalah', slide_num=2)
    add_footer_bar(slide)

    items = [
        'SMKN 1 Marabahan memiliki 575 siswa aktif (per April 2026), '
        'proses presensi masih menggunakan buku kertas — rekap lambat, '
        'rawan manipulasi data.',
        'Fingerprint juga kurang efektif: sering gagal baca sidik jari '
        '(basah/kotor), antrean panjang, kontak fisik tidak higienis.',
        'Solusi: sistem presensi berbasis Face Recognition via Android — '
        'siswa cukup menatap kamera, proses cepat & tanpa kontak fisik.',
        'Dilengkapi Liveness Detection (deteksi kedipan mata) untuk '
        'mencegah kecurangan menggunakan foto/video (anti-spoofing).',
        'Sistem terintegrasi dengan web dashboard untuk monitoring '
        'oleh Guru, TU, dan Orang Tua siswa.',
    ]
    # Problem icon area
    add_bg_rect(slide, Inches(0.4), Inches(1.5), Inches(0.06), Inches(4.5), BLUE_PRIMARY)

    add_bullet_text(slide, CONTENT_MARGIN_L, CONTENT_TOP,
                    CONTENT_W, Inches(5.0), items, font_size=16, color=BLACK)

    set_morph_transition(slide)
    return slide


def create_slide3(pres):
    """Slide 3: Rumusan Masalah & Batasan Masalah"""
    slide = pres.slides.add_slide(pres.slide_layouts[6])
    add_section_header(slide, 'Rumusan Masalah & Batasan Masalah', slide_num=3)
    add_footer_bar(slide)

    # Rumusan
    add_bg_rect(slide, Inches(0.4), Inches(1.5), Inches(0.06), Inches(1.0), BLUE_PRIMARY)

    # Rumusan label
    add_textbox(slide, CONTENT_MARGIN_L, Inches(1.5), Inches(2.0), Inches(0.4),
                'RUMUSAN MASALAH', font_size=16, bold=True, color=BLUE_PRIMARY,
                font_name='Calibri')

    add_textbox(slide, CONTENT_MARGIN_L, Inches(1.95), CONTENT_W, Inches(0.6),
                '"Bagaimana mengimplementasikan teknologi pengenalan wajah (face recognition) '
                'pada aplikasi presensi Android guna memfasilitasi input data kehadiran '
                'secara efisien di SMKN 1 Marabahan?"',
                font_size=14, color=BLACK, font_name='Calibri')

    # Batasan
    add_bg_rect(slide, Inches(0.4), Inches(2.8), Inches(0.06), Inches(4.0), CYAN)

    add_textbox(slide, CONTENT_MARGIN_L, Inches(2.8), Inches(3.0), Inches(0.4),
                'BATASAN MASALAH', font_size=16, bold=True, color=CYAN,
                font_name='Calibri')

    items = [
        'Lingkungan: SMKN 1 Marabahan',
        '4 peran pengguna: Siswa (Android), Orang Tua (Web), Guru (Web), TU (Web)',
        'Metode Waterfall (analisis hingga pengujian)',
        'UML: Use Case Diagram, Class Diagram, ERD',
        'Teknologi: Kotlin (Android), Python/FastAPI (Backend), React JS (Web)',
        'Pengujian: Black-box Testing',
        'Pre-trained model ArcFace + RetinaFace (DeepFace)',
    ]
    add_bullet_text(slide, CONTENT_MARGIN_L, Inches(3.25),
                    CONTENT_W, Inches(3.5), items, font_size=14, color=BLACK,
                    bullet_char='\u25B8', line_spacing=1.4)

    set_morph_transition(slide)
    return slide


def create_slide4(pres):
    """Slide 4: Tujuan & Manfaat"""
    slide = pres.slides.add_slide(pres.slide_layouts[6])
    add_section_header(slide, 'Tujuan & Manfaat Penelitian', slide_num=4)
    add_footer_bar(slide)

    # Tujuan
    add_bg_rect(slide, Inches(0.4), Inches(1.5), Inches(0.06), Inches(1.0), BLUE_PRIMARY)
    add_textbox(slide, CONTENT_MARGIN_L, Inches(1.5), Inches(2.0), Inches(0.4),
                'TUJUAN', font_size=16, bold=True, color=BLUE_PRIMARY, font_name='Calibri')
    add_textbox(slide, CONTENT_MARGIN_L, Inches(1.95), CONTENT_W, Inches(0.5),
                'Mengimplementasikan face recognition pada aplikasi presensi Android '
                'untuk input data kehadiran yang efisien di SMKN 1 Marabahan.',
                font_size=14, color=BLACK, font_name='Calibri')

    # Manfaat - split into 4 columns
    benefits = [
        ('Sekolah', GREEN, 'Transformasi digital,\nkedisiplinan siswa,\n efisiensi operasional'),
        ('Orang Tua', BLUE_PRIMARY, 'Pantau kehadiran\nanak secara\nreal-time via web'),
        ('Siswa', CYAN, 'Pengalaman AI,\npresensi cepat\ntanpa antre'),
        ('Peneliti', RGBColor(0xF5, 0x9E, 0x0B), 'Implementasi\nRetinaFace +\nArcFace'),
    ]

    col_w = Inches(2.6)
    col_gap = Inches(0.3)
    total_w = col_w * 4 + col_gap * 3
    start_x = (SLIDE_W - total_w) // 2

    add_textbox(slide, CONTENT_MARGIN_L, Inches(2.7), Inches(3.0), Inches(0.4),
                'MANFAAT', font_size=16, bold=True, color=BLUE_PRIMARY, font_name='Calibri')

    for i, (title, accent_color, desc) in enumerate(benefits):
        x = start_x + (col_w + col_gap) * i
        y = Inches(3.3)
        # Card background
        add_bg_rect(slide, x, y, col_w, Inches(2.8), GRAY_LIGHT)
        # Accent top line
        add_bg_rect(slide, x, y, col_w, Inches(0.06), accent_color)
        # Card title
        add_textbox(slide, x + Inches(0.2), y + Inches(0.2), col_w - Inches(0.4), Inches(0.5),
                    title, font_size=18, bold=True, color=accent_color,
                    alignment=PP_ALIGN.CENTER, font_name='Calibri')
        # Card description
        add_textbox(slide, x + Inches(0.2), y + Inches(0.8), col_w - Inches(0.4), Inches(1.8),
                    desc, font_size=13, color=GRAY, alignment=PP_ALIGN.CENTER,
                    font_name='Calibri')

    set_morph_transition(slide)
    return slide


def create_slide5(pres):
    """Slide 5: Dasar Teori"""
    slide = pres.slides.add_slide(pres.slide_layouts[6])
    add_section_header(slide, 'Dasar Teori', slide_num=5)
    add_footer_bar(slide)

    # Tech stack cards - 4 columns
    techs = [
        ('Face Recognition', BLUE_PRIMARY, [
            'ArcFace: loss function dgn\nAdditive Angular Margin',
            'RetinaFace: one-stage\nface detector (2019)',
            'DeepFace library dgn\nCosine Similarity (>= 0.60)',
        ]),
        ('Android App', CYAN, [
            'Kotlin + CameraX 1.4.1',
            'MLKit Face Detection\n(pose yaw angle)',
            'Retrofit 2.9 + OkHttp\n(multipart upload)',
        ]),
        ('Backend API', GREEN, [
            'Python FastAPI',
            'SQLAlchemy ORM + MySQL',
            'JWT Authentication\n(HS256, 24h expiry)',
        ]),
        ('Web Dashboard', RGBColor(0xF5, 0x9E, 0x0B), [
            'React 19 + Vite 8',
            'Ant Design 6 UI',
            'Axios + Chart Analytics',
        ]),
    ]

    col_w = Inches(2.85)
    col_gap = Inches(0.2)
    total_w = col_w * 4 + col_gap * 3
    start_x = (SLIDE_W - total_w) // 2

    for i, (title, accent_color, items) in enumerate(techs):
        x = start_x + (col_w + col_gap) * i
        y = Inches(1.6)
        # Card bg
        add_bg_rect(slide, x, y, col_w, Inches(3.2), GRAY_LIGHT)
        add_bg_rect(slide, x, y, col_w, Inches(0.06), accent_color)
        # Title
        add_textbox(slide, x + Inches(0.15), y + Inches(0.2), col_w - Inches(0.3), Inches(0.5),
                    title, font_size=17, bold=True, color=accent_color,
                    alignment=PP_ALIGN.CENTER, font_name='Calibri')
        # Items
        for j, item in enumerate(items):
            add_textbox(slide, x + Inches(0.15), y + Inches(0.85) + Inches(0.75) * j,
                        col_w - Inches(0.3), Inches(0.7),
                        f'\u2022 {item}', font_size=12, color=BLACK, font_name='Calibri')

    # Bottom: Methodology
    add_bg_rect(slide, CONTENT_MARGIN_L, Inches(5.2), CONTENT_W, Inches(1.4), GRAY_LIGHT)
    add_bg_rect(slide, CONTENT_MARGIN_L, Inches(5.2), CONTENT_W, Inches(0.06), BLUE_PRIMARY)
    add_textbox(slide, CONTENT_MARGIN_L + Inches(0.2), Inches(5.35), Inches(3.0), Inches(0.4),
                'Metodologi & Tools', font_size=14, bold=True, color=BLUE_PRIMARY, font_name='Calibri')
    add_textbox(slide, CONTENT_MARGIN_L + Inches(0.2), Inches(5.75), CONTENT_W - Inches(0.4), Inches(0.7),
                'Metode Waterfall  |  UML (Use Case, Class Diagram, ERD)  |  '
                'Figma (Wireframe)  |  VS Code + Android Studio  |  Black-box Testing',
                font_size=13, color=GRAY, font_name='Calibri')

    set_morph_transition(slide)
    return slide


def create_slide6(pres):
    """Slide 6: Arsitektur Sistem"""
    slide = pres.slides.add_slide(pres.slide_layouts[6])
    add_section_header(slide, 'Arsitektur Sistem', slide_num=6)
    add_footer_bar(slide)

    # Architecture blocks
    blocks = [
        ('Aplikasi Android\n(Kotlin)', Inches(0.5), Inches(2.0), Inches(3.0), Inches(2.5), BLUE_PRIMARY, [
            'CameraX + MLKit',
            'Face Detection',
            '3 Pose Capture',
            'Retrofit HTTP Client',
        ]),
        ('FastAPI Server\n(Python)', Inches(4.8), Inches(1.6), Inches(3.5), Inches(3.3), CYAN, [
            'Face Registration',
            'Face Attendance',
            'Auth (JWT)',
            'CRUD API',
            'DeepFace ArcFace',
        ]),
        ('Database\nMySQL', Inches(4.8), Inches(5.2), Inches(3.5), Inches(1.2), RGBColor(0xF5, 0x9E, 0x0B), [
            '8 Tabel: jurusan, kelas, siswa,\nface_data, absensi, users, guru, guru_kelas',
        ]),
        ('Web Dashboard\n(React JS)', Inches(9.6), Inches(2.0), Inches(3.0), Inches(2.5), GREEN, [
            'Ant Design UI',
            'Login (Guru/TU/Ortu)',
            'Dashboard Analytics',
            'CRUD & Export Excel',
        ]),
    ]

    for caption, x, y, w, h, color, details in blocks:
        # Block background
        add_bg_rect(slide, x, y, w, h, color)
        # Block title
        add_textbox(slide, x + Inches(0.15), y + Inches(0.15), w - Inches(0.3), Inches(0.6),
                    caption, font_size=14, bold=True, color=WHITE,
                    alignment=PP_ALIGN.CENTER, font_name='Calibri')
        # Divider line
        add_bg_rect(slide, x + Inches(0.3), y + Inches(0.8), w - Inches(0.6), Inches(0.02), RGBColor(0xFF, 0xFF, 0xFF))
        # Details
        detail_text = '\n'.join(f'\u2022 {d}' for d in details)
        add_textbox(slide, x + Inches(0.15), y + Inches(0.9), w - Inches(0.3), h - Inches(1.0),
                    detail_text, font_size=11, color=WHITE, font_name='Calibri')

    # Arrows / connections
    arrow_texts = [
        (Inches(3.5), Inches(3.0), 'Retrofit\nHTTP'),
        (Inches(8.3), Inches(3.0), 'Axios\nHTTP'),
        (Inches(8.3), Inches(5.2), 'SQLAlchemy\nORM'),
    ]
    for ax, ay, label in arrow_texts:
        add_textbox(slide, ax, ay, Inches(1.3), Inches(0.6),
                    label, font_size=11, bold=True, color=GRAY,
                    alignment=PP_ALIGN.CENTER, font_name='Calibri')

    # Face pipeline note
    add_bg_rect(slide, Inches(0.5), Inches(5.2), Inches(3.0), Inches(1.2), GRAY_LIGHT)
    add_textbox(slide, Inches(0.65), Inches(5.3), Inches(2.7), Inches(1.0),
                'Face Pipeline:\nCamera \u2192 MLKit Detect \u2192\nDeepFace ArcFace \u2192\nCosine Similarity \u2265 0.60',
                font_size=11, color=BLACK, font_name='Calibri')

    # Legend
    add_textbox(slide, Inches(10), Inches(5.5), Inches(2.8), Inches(0.8),
                'Alur Data:\nAndroid \u2194 API \u2194 DB\nWeb \u2194 API \u2194 DB',
                font_size=10, color=GRAY, font_name='Calibri')

    set_morph_transition(slide)
    return slide


def create_slide7(pres):
    """Slide 7: Metode Penelitian (1) - Analisis & Perancangan"""
    slide = pres.slides.add_slide(pres.slide_layouts[6])
    add_section_header(slide, 'Metode Waterfall \u2013 Tahap 1 & 2',
                       'Analisis Kebutuhan & Perancangan Sistem', slide_num=7)
    add_footer_bar(slide)

    # Tahap 1
    add_bg_rect(slide, Inches(0.5), Inches(1.6), Inches(5.8), Inches(2.6), GRAY_LIGHT)
    add_bg_rect(slide, Inches(0.5), Inches(1.6), Inches(0.06), Inches(2.6), BLUE_PRIMARY)
    add_textbox(slide, Inches(0.8), Inches(1.7), Inches(5.0), Inches(0.4),
                'Tahap 1 : Analisis Kebutuhan', font_size=18, bold=True, color=BLUE_PRIMARY, font_name='Calibri')

    analisis_items = [
        'Observasi langsung di SMKN 1 Marabahan\n(4 Mei 2026, proses presesi pagi hari)',
        'Wawancara dengan Bapak Ridwan (Staff TU)\nmengenai alur presensi & kendala',
        'Studi pustaka: 9+ jurnal referensi',
        'Identifikasi kebutuhan fungsional\n& non-fungsional sistem',
    ]
    add_bullet_text(slide, Inches(0.8), Inches(2.2), Inches(5.2), Inches(2.0),
                    analisis_items, font_size=13, color=BLACK, bullet_char='\u25B8', line_spacing=1.5)

    # Tahap 2
    add_bg_rect(slide, Inches(6.7), Inches(1.6), Inches(5.8), Inches(2.6), GRAY_LIGHT)
    add_bg_rect(slide, Inches(6.7), Inches(1.6), Inches(0.06), Inches(2.6), CYAN)
    add_textbox(slide, Inches(7.0), Inches(1.7), Inches(5.0), Inches(0.4),
                'Tahap 2 : Perancangan Sistem', font_size=18, bold=True, color=CYAN, font_name='Calibri')

    design_items = [
        'Arsitektur client-server\n(Android + Web + API)',
        'Use Case Diagram (4 aktor)',
        'Class Diagram (8 kelas model)',
        'Entity Relationship Diagram (ERD)',
        'Wireframe medium fidelity (Figma)',
    ]
    add_bullet_text(slide, Inches(7.0), Inches(2.2), Inches(5.2), Inches(2.0),
                    design_items, font_size=13, color=BLACK, bullet_char='\u25B8', line_spacing=1.5)

    # Bottom info
    add_bg_rect(slide, Inches(0.5), Inches(4.6), Inches(12.0), Inches(1.8), GRAY_LIGHT)
    add_bg_rect(slide, Inches(0.5), Inches(4.6), Inches(12.0), Inches(0.06), BLUE_PRIMARY)
    add_textbox(slide, Inches(0.8), Inches(4.8), Inches(11.4), Inches(0.4),
                'Hasil Perancangan:', font_size=14, bold=True, color=BLUE_PRIMARY, font_name='Calibri')

    results_items = [
        'Sistem Android untuk presensi siswa dengan fitur face recognition (3 pose capture + scan wajah)',
        'Web Dashboard untuk Guru (rekap absensi, ubah status) dan TU (CRUD master, export Excel)',
        'Web Portal untuk Orang Tua (pantau kehadiran anak via NISN)',
        'Backend API FastAPI dengan DeepFace ArcFace untuk verifikasi wajah (cosine similarity \u2265 0.60)',
    ]
    add_bullet_text(slide, Inches(0.8), Inches(5.25), Inches(11.4), Inches(1.2),
                    results_items, font_size=12, color=BLACK, bullet_char='\u25B8', line_spacing=1.4)

    set_morph_transition(slide)
    return slide


def create_slide8(pres):
    """Slide 8: Metode Penelitian (2) - Implementasi & Pengujian"""
    slide = pres.slides.add_slide(pres.slide_layouts[6])
    add_section_header(slide, 'Metode Waterfall \u2013 Tahap 3 & 4',
                       'Implementasi & Pengujian Sistem', slide_num=8)
    add_footer_bar(slide)

    # Tahap 3
    add_bg_rect(slide, Inches(0.5), Inches(1.6), Inches(5.8), Inches(3.0), GRAY_LIGHT)
    add_bg_rect(slide, Inches(0.5), Inches(1.6), Inches(0.06), Inches(3.0), GREEN)
    add_textbox(slide, Inches(0.8), Inches(1.7), Inches(5.0), Inches(0.4),
                'Tahap 3 : Implementasi', font_size=18, bold=True, color=GREEN, font_name='Calibri')

    impl_items = [
        'Android Studio (Kotlin + XML) \u2014\nAplikasi presensi dengan CameraX &\nMLKit Face Detection',
        'VS Code (Python FastAPI) \u2014\nBackend REST API dengan SQLAlchemy\nORM & MySQL',
        'VS Code (React JS + Ant Design) \u2014\nWeb Dashboard responsive',
        'DeepFace (ArcFace + RetinaFace) \u2014\nFace recognition engine',
    ]
    add_bullet_text(slide, Inches(0.8), Inches(2.2), Inches(5.2), Inches(2.3),
                    impl_items, font_size=13, color=BLACK, bullet_char='\u25B8', line_spacing=1.5)

    # Tahap 4
    add_bg_rect(slide, Inches(6.7), Inches(1.6), Inches(5.8), Inches(3.0), GRAY_LIGHT)
    add_bg_rect(slide, Inches(6.7), Inches(1.6), Inches(0.06), Inches(3.0), RGBColor(0xF5, 0x9E, 0x0B))
    add_textbox(slide, Inches(7.0), Inches(1.7), Inches(5.0), Inches(0.4),
                'Tahap 4 : Pengujian', font_size=18, bold=True,
                color=RGBColor(0xF5, 0x9E, 0x0B), font_name='Calibri')

    test_text = (
        'Metode: Black-box Testing\n'
        'Fokus: Fungsionalitas fitur\n'
        'Skenario uji:\n'
        '\u25B8 Registrasi wajah (3 pose)\n'
        '\u25B8 Scan presensi (Hadir/Tidak Dikenal)\n'
        '\u25B8 CRUD siswa & guru\n'
        '\u25B8 Login role-based (Guru/TU/Siswa)\n'
        '\u25B8 Export Excel absensi'
    )
    add_textbox(slide, Inches(7.0), Inches(2.2), Inches(5.2), Inches(2.3),
                test_text, font_size=13, color=BLACK, font_name='Calibri')

    # Toll section at bottom
    add_bg_rect(slide, Inches(0.5), Inches(5.0), Inches(12.0), Inches(1.5), GRAY_LIGHT)
    add_bg_rect(slide, Inches(0.5), Inches(5.0), Inches(12.0), Inches(0.06), BLUE_PRIMARY)
    add_textbox(slide, Inches(0.8), Inches(5.2), Inches(3.0), Inches(0.4),
                'Lingkungan Pengembangan:', font_size=14, bold=True, color=BLUE_PRIMARY, font_name='Calibri')

    tools_text = ('Android Studio (IDE Android)  |  VS Code (IDE Backend & Web)  |  '
                  'Laragon (Local Server)  |  Figma (Wireframe)  |  '
                  'MySQL (Database)  |  Git (Version Control)')
    add_textbox(slide, Inches(0.8), Inches(5.7), Inches(11.4), Inches(0.6),
                tools_text, font_size=13, color=GRAY, font_name='Calibri')

    set_morph_transition(slide)
    return slide


def create_slide9(pres):
    """Slide 9: Diagram Perancangan"""
    slide = pres.slides.add_slide(pres.slide_layouts[6])
    add_section_header(slide, 'Diagram Perancangan Sistem',
                       'Use Case Diagram  |  Class Diagram  |  ERD', slide_num=9)
    add_footer_bar(slide)

    # 3 columns for 3 diagrams
    diagrams = [
        ('Use Case Diagram', BLUE_PRIMARY, [
            '4 aktor:',
            '\u25B8 Siswa (Android)',
            '\u25B8 Orang Tua (Web)',
            '\u25B8 Guru (Web)',
            '\u25B8 TU (Web)',
            'Fitur utama:',
            '\u25B8 Registrasi Wajah',
            '\u25B8 Presensi Face Scan',
            '\u25B8 Monitoring Ortu',
            '\u25B8 Dashboard & CRUD',
        ]),
        ('Class Diagram', CYAN, [
            '8 kelas model:',
            '\u25B8 Jurusan, Kelas, Siswa',
            '\u25B8 FaceData, Absensi',
            '\u25B8 User, Guru, GuruKelas',
            'Android:',
            '\u25B8 ApiService, ApiClient',
            '\u25B8 MainActivity',
            '\u25B8 RegisterActivity',
            '\u25B8 CameraActivity',
        ]),
        ('ERD', GREEN, [
            'Relasi database:',
            '\u25B8 Jurusan 1\u2013* Kelas',
            '\u25B8 Kelas 1\u2013* Siswa',
            '\u25B8 Siswa 1\u2013* FaceData',
            '\u25B8 Siswa 1\u2013* Absensi',
            '\u25B8 User 1\u20131 Guru',
            '\u25B8 User *\u2013* Kelas',
            '\u25B8 User 1\u2013* Absensi',
        ]),
    ]

    col_w = Inches(3.8)
    col_gap = Inches(0.3)
    total_w = col_w * 3 + col_gap * 2
    start_x = (SLIDE_W - total_w) // 2

    for i, (title, accent_color, items) in enumerate(diagrams):
        x = start_x + (col_w + col_gap) * i
        y = Inches(1.6)
        # Card
        add_bg_rect(slide, x, y, col_w, Inches(4.5), GRAY_LIGHT)
        add_bg_rect(slide, x, y, col_w, Inches(0.06), accent_color)
        # Title
        add_textbox(slide, x + Inches(0.15), y + Inches(0.2), col_w - Inches(0.3), Inches(0.5),
                    title, font_size=18, bold=True, color=accent_color,
                    alignment=PP_ALIGN.CENTER, font_name='Calibri')
        # Separator
        add_bg_rect(slide, x + Inches(0.5), y + Inches(0.7), col_w - Inches(1.0), Inches(0.02), accent_color)
        # Items
        item_text = '\n'.join(items)
        add_textbox(slide, x + Inches(0.2), y + Inches(0.85), col_w - Inches(0.4), Inches(3.5),
                    item_text, font_size=12, color=BLACK, font_name='Calibri')

    # Note at bottom
    add_textbox(slide, Inches(1.5), Inches(6.4), Inches(10.3), Inches(0.5),
                'Diagram detail tersedia di file PlantUML: docs/use-case-diagram.puml, docs/class-diagram.puml, docs/erd.puml',
                font_size=11, color=GRAY, alignment=PP_ALIGN.CENTER, font_name='Calibri')

    set_morph_transition(slide)
    return slide


def create_slide10(pres):
    """Slide 10: Penutup"""
    slide = pres.slides.add_slide(pres.slide_layouts[6])

    # Blue background
    add_bg_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, BLUE_PRIMARY)

    # Logo
    if os.path.exists(LOGO_PATH):
        img = PILImage.open(LOGO_PATH)
        w, h = img.size
        scale = 1.5 / (w / 96)
        logo_w = Inches(1.5)
        logo_h = Inches(h / 96 * scale)
        logo_left = (SLIDE_W - logo_w) // 2
        slide.shapes.add_picture(LOGO_PATH, logo_left, Inches(0.8), logo_w, logo_h)

    # Thank you text
    add_textbox(slide, Inches(0), Inches(2.5), SLIDE_W, Inches(1.5),
                'Terima Kasih', font_size=48, bold=True, color=WHITE,
                alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name='Calibri')

    # Q&A
    add_textbox(slide, Inches(0), Inches(3.8), SLIDE_W, Inches(1.0),
                'Sesi Tanya Jawab', font_size=28, color=RGBColor(0xCC, 0xD5, 0xFF),
                alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name='Calibri')

    # Divider line
    add_bg_rect(slide, Inches(5), Inches(4.8), Inches(3.3), Inches(0.03), WHITE)

    # Contact info
    contact_text = 'Ahmad Nuur Fu\'ady  |  24302064\nProgram Studi D3 Teknik Informatika\nPoliteknik Hasnur'
    add_textbox(slide, Inches(0), Inches(5.1), SLIDE_W, Inches(1.2),
                contact_text, font_size=16, color=RGBColor(0xAA, 0xBB, 0xFF),
                alignment=PP_ALIGN.CENTER, font_name='Calibri')

    add_slide_number(slide, 10)
    set_morph_transition(slide)
    return slide


# ─── MAIN ────────────────────────────────────────────────────────

def main():
    pres = Presentation()
    pres.slide_width = SLIDE_W
    pres.slide_height = SLIDE_H

    create_slide1(pres)
    create_slide2(pres)
    create_slide3(pres)
    create_slide4(pres)
    create_slide5(pres)
    create_slide6(pres)
    create_slide7(pres)
    create_slide8(pres)
    create_slide9(pres)
    create_slide10(pres)

    pres.save(OUTPUT_PATH)
    print(f'Presentation saved to: {OUTPUT_PATH}')
    print(f'Total slides: {len(pres.slides)}')

if __name__ == '__main__':
    main()
